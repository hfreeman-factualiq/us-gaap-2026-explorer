#!/usr/bin/env python3
"""Cached, parallel text extraction for the company document packs.

Parsing PDFs, decks, memos and workbooks is the CPU-heavy half of the company
tree build, so extraction runs in a process pool and every result is cached on
disk keyed by file size + mtime. Reruns read the cache instead of the documents.
"""

from __future__ import annotations

import hashlib
import json
import logging
import warnings
from concurrent.futures import ProcessPoolExecutor, as_completed
from pathlib import Path

ROOT = Path(__file__).resolve().parent
CACHE_ROOT = ROOT / ".cache"
TEXT_CACHE = CACHE_ROOT / "doc_text"
JSON_CACHE = CACHE_ROOT / "doc_json"
CACHE_VERSION = "1"

log = logging.getLogger("extract")

NARRATIVE_SUFFIXES = {".pdf", ".docx", ".pptx", ".msg", ".txt"}

# Parsers complain loudly about real-world files (broken xref tables, bad number
# formats). The complaints are benign and drown out progress output.
_NOISY_LOGGERS = (
    "pypdf",
    "pypdf.generic",
    "pypdf._reader",
    "pypdf._page",
    "pypdf.filters",
    "extract_msg",
    "pptx",
    "openpyxl",
)


def quiet_parser_noise() -> None:
    for name in _NOISY_LOGGERS:
        logging.getLogger(name).setLevel(logging.CRITICAL)
    warnings.filterwarnings("ignore")


def fingerprint(path: Path, extra: str = "") -> str:
    st = path.stat()
    raw = f"{path.resolve()}|{st.st_size}|{st.st_mtime_ns}|{CACHE_VERSION}|{extra}"
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:32]


def json_cache_get(key: str):
    path = JSON_CACHE / f"{key}.json"
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:  # noqa: BLE001 - a corrupt cache entry is just a miss
        return None


def json_cache_put(key: str, payload) -> None:
    JSON_CACHE.mkdir(parents=True, exist_ok=True)
    (JSON_CACHE / f"{key}.json").write_text(
        json.dumps(payload, ensure_ascii=False), encoding="utf-8"
    )


# --- format-specific text extraction ---------------------------------------


def extract_pdf_text(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def extract_docx_text(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    bits = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            bits.append(" | ".join(cell.text for cell in row.cells if cell.text))
    return "\n".join(bits)


def extract_pptx_text(path: Path) -> str:
    from pptx import Presentation

    deck = Presentation(str(path))
    bits: list[str] = []
    for slide in deck.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                bits.append(text)
    return "\n".join(bits)


def extract_msg_text(path: Path) -> str:
    import extract_msg

    msg = extract_msg.Message(str(path))
    try:
        return "\n".join(str(x) for x in [msg.subject, msg.body] if x)
    finally:
        msg.close()


def extract_narrative_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path)
    if suffix == ".docx":
        return extract_docx_text(path)
    if suffix == ".pptx":
        return extract_pptx_text(path)
    if suffix == ".msg":
        return extract_msg_text(path)
    if suffix == ".txt":
        return path.read_text(encoding="utf-8", errors="replace")
    return ""


# --- cached / parallel driver ----------------------------------------------


def _text_worker(path_str: str) -> dict:
    quiet_parser_noise()
    path = Path(path_str)
    result = {"path": path_str, "cached": False, "error": None, "chars": 0}
    try:
        key = fingerprint(path)
    except OSError as exc:
        result["error"] = str(exc)[:300]
        return result
    cache_file = TEXT_CACHE / f"{key}.txt"
    if cache_file.exists():
        result["cached"] = True
        result["cache_file"] = str(cache_file)
        result["chars"] = cache_file.stat().st_size
        return result
    try:
        text = extract_narrative_text(path)
    except Exception as exc:  # noqa: BLE001 - one bad document must not stop the run
        result["error"] = f"{type(exc).__name__}: {exc}"[:300]
        return result
    # PDF text layers can yield lone surrogates, which are not encodable.
    text = text.encode("utf-8", "replace").decode("utf-8", "replace")
    TEXT_CACHE.mkdir(parents=True, exist_ok=True)
    tmp = cache_file.with_suffix(".tmp")
    tmp.write_text(text, encoding="utf-8")
    tmp.replace(cache_file)
    result["cache_file"] = str(cache_file)
    result["chars"] = len(text)
    return result


def read_cached_text(path: Path) -> str:
    cache_file = TEXT_CACHE / f"{fingerprint(path)}.txt"
    if not cache_file.exists():
        return ""
    return cache_file.read_text(encoding="utf-8", errors="replace")


def extract_texts(
    paths: list[Path],
    *,
    max_workers: int,
    on_done=None,
) -> tuple[dict[str, str], list[dict]]:
    """Extract text for every path in parallel; return (path→text, errors)."""
    texts: dict[str, str] = {}
    errors: list[dict] = []
    if not paths:
        return texts, errors

    unique = list(dict.fromkeys(str(p) for p in paths))
    workers = max(1, min(max_workers, len(unique)))
    done = 0
    with ProcessPoolExecutor(max_workers=workers) as pool:
        futures = {pool.submit(_text_worker, p): p for p in unique}
        for fut in as_completed(futures):
            path_str = futures[fut]
            done += 1
            try:
                res = fut.result()
            except Exception as exc:  # noqa: BLE001 - worker crash (e.g. OOM)
                res = {"path": path_str, "error": f"worker: {exc}"[:300], "cached": False}
            if res.get("error"):
                errors.append({"file": path_str, "error": res["error"]})
            elif res.get("cache_file"):
                texts[path_str] = Path(res["cache_file"]).read_text(
                    encoding="utf-8", errors="replace"
                )
            if on_done:
                on_done(done, len(unique), res)
    return texts, errors
